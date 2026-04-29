from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BLUE_DARK = RGBColor(15, 45, 85)
BLUE = RGBColor(29, 119, 229)
BLUE_LIGHT = RGBColor(226, 241, 255)
GRAY_BG = RGBColor(246, 248, 251)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(107, 114, 128)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(20, 184, 166)
ORANGE = RGBColor(245, 158, 11)


class PPTMaster:
    def __init__(self, spec: Dict[str, Any]):
        self.spec = spec
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.w = self.prs.slide_width
        self.h = self.prs.slide_height

    def build(self) -> Presentation:
        self.add_cover()
        self.add_agenda()
        for section in self.spec.get("sections", []):
            self.add_section(section)
            for slide in section.get("slides", []):
                self.add_content_slide(slide)
        self.add_thanks()
        return self.prs

    def blank(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = GRAY_BG
        return slide

    def add_top_bar(self, slide, label: str = ""):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.w, Inches(0.13))
        shape.fill.solid(); shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        if label:
            tx = slide.shapes.add_textbox(Inches(11.1), Inches(0.22), Inches(1.8), Inches(0.25))
            p = tx.text_frame.paragraphs[0]
            p.text = label
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(9); p.font.color.rgb = MUTED

    def add_text(self, slide, text, x, y, w, h, size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        return box

    def add_card(self, slide, x, y, w, h, title: str = "", body: str = "", accent=BLUE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(225, 232, 240)
        card.line.width = Pt(1)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.09), h)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        if title:
            self.add_text(slide, title, x+Inches(0.25), y+Inches(0.16), w-Inches(0.45), Inches(0.36), 17, True, BLUE_DARK)
        if body:
            box = slide.shapes.add_textbox(x+Inches(0.25), y+Inches(0.58), w-Inches(0.45), h-Inches(0.7))
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            for i, line in enumerate(body.split("\n")):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(12.5)
                p.font.color.rgb = TEXT
                p.space_after = Pt(4)

    def add_cover(self):
        slide = self.blank()
        self.add_top_bar(slide)
        school = self.spec.get("school", "")
        if school:
            self.add_text(slide, school, Inches(0.72), Inches(0.5), Inches(5), Inches(0.35), 15, False, MUTED)
        self.add_text(slide, self.spec.get("title", "Untitled Presentation"), Inches(0.72), Inches(1.5), Inches(9.8), Inches(1.0), 34, True, BLUE_DARK)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(2.75), Inches(2.2), Inches(0.05))
        line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.fill.background()
        subtitle = self.spec.get("subtitle", "")
        if subtitle:
            self.add_text(slide, subtitle, Inches(0.72), Inches(3.05), Inches(8.5), Inches(0.5), 19, False, TEXT)
        meta = self.spec.get("meta", [])
        self.add_card(slide, Inches(0.72), Inches(4.25), Inches(5.6), Inches(1.55), "答辩信息 / Project Info", "\n".join(meta), BLUE)
        # tech visual
        for i in range(5):
            x = Inches(9.0 + i*0.45)
            y = Inches(1.2 + i*0.52)
            s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.16+i*0.05), Inches(0.16+i*0.05))
            s.fill.solid(); s.fill.fore_color.rgb = RGBColor(120, 200, 255)
            s.line.fill.background()
        self.add_text(slide, "AI Agent · Automation · Engineering", Inches(8.0), Inches(5.85), Inches(4.5), Inches(0.4), 13, False, MUTED, PP_ALIGN.RIGHT)

    def add_agenda(self):
        slide = self.blank(); self.add_top_bar(slide, "Agenda")
        self.add_text(slide, "目录 / Agenda", Inches(0.7), Inches(0.55), Inches(5), Inches(0.45), 26, True, BLUE_DARK)
        sections = self.spec.get("sections", [])
        for i, sec in enumerate(sections[:6]):
            x = Inches(0.85 + (i % 2) * 6.05)
            y = Inches(1.45 + (i // 2) * 1.55)
            num = f"0{i+1}"
            self.add_card(slide, x, y, Inches(5.55), Inches(1.1), sec.get("title", "Section"), sec.get("desc", ""), [BLUE, GREEN, ORANGE][i % 3])
            self.add_text(slide, num, x+Inches(4.75), y+Inches(0.1), Inches(0.55), Inches(0.25), 12, True, MUTED, PP_ALIGN.RIGHT)

    def add_section(self, section: Dict[str, Any]):
        slide = self.blank(); self.add_top_bar(slide, "Section")
        self.add_text(slide, section.get("title", "Section"), Inches(0.75), Inches(2.35), Inches(8.8), Inches(0.75), 32, True, BLUE_DARK)
        desc = section.get("desc", "")
        if desc:
            self.add_text(slide, desc, Inches(0.78), Inches(3.18), Inches(8.6), Inches(0.45), 17, False, TEXT)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(3.9), Inches(2.8), Inches(0.05))
        line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.fill.background()

    def add_content_slide(self, data: Dict[str, Any]):
        slide = self.blank(); self.add_top_bar(slide, data.get("tag", ""))
        self.add_text(slide, data.get("title", "Slide"), Inches(0.65), Inches(0.42), Inches(10.5), Inches(0.45), 24, True, BLUE_DARK)
        layout = data.get("layout", "cards")
        if layout == "two_col":
            items = data.get("items", [])
            self.add_card(slide, Inches(0.75), Inches(1.35), Inches(5.75), Inches(4.85), items[0].get("title", ""), items[0].get("body", "") if items else "", BLUE)
            if len(items) > 1:
                self.add_card(slide, Inches(6.8), Inches(1.35), Inches(5.75), Inches(4.85), items[1].get("title", ""), items[1].get("body", ""), GREEN)
        elif layout == "timeline":
            self.draw_timeline(slide, data.get("items", []))
        elif layout == "metrics":
            self.draw_metrics(slide, data.get("items", []))
        elif layout == "flow":
            self.draw_flow(slide, data.get("items", []))
        else:
            items = data.get("items", [])
            for i, item in enumerate(items[:6]):
                x = Inches(0.75 + (i % 3) * 4.15)
                y = Inches(1.25 + (i // 3) * 2.25)
                self.add_card(slide, x, y, Inches(3.8), Inches(1.75), item.get("title", ""), item.get("body", ""), [BLUE, GREEN, ORANGE][i % 3])
        note = data.get("note")
        if note:
            self.add_text(slide, note, Inches(0.8), Inches(6.85), Inches(11.8), Inches(0.25), 10, False, MUTED)

    def draw_metrics(self, slide, items: List[Dict[str, Any]]):
        for i, item in enumerate(items[:4]):
            x = Inches(0.82 + i * 3.05)
            y = Inches(2.2)
            self.add_card(slide, x, y, Inches(2.72), Inches(2.25), item.get("label", item.get("title", "")), item.get("body", ""), [BLUE, GREEN, ORANGE, BLUE_DARK][i%4])
            self.add_text(slide, item.get("value", "--"), x+Inches(0.25), y+Inches(0.78), Inches(2.2), Inches(0.5), 28, True, BLUE_DARK, PP_ALIGN.CENTER)

    def draw_flow(self, slide, items: List[Dict[str, Any]]):
        for i, item in enumerate(items[:5]):
            x = Inches(0.75 + i * 2.45)
            y = Inches(2.35)
            self.add_card(slide, x, y, Inches(2.0), Inches(1.35), item.get("title", f"Step {i+1}"), item.get("body", ""), BLUE)
            if i < min(len(items), 5)-1:
                self.add_text(slide, "→", x+Inches(2.05), y+Inches(0.4), Inches(0.25), Inches(0.3), 22, True, MUTED)

    def draw_timeline(self, slide, items: List[Dict[str, Any]]):
        base_y = Inches(3.2)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), base_y, Inches(11.2), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.fill.background()
        for i, item in enumerate(items[:5]):
            x = Inches(1.1 + i * 2.55)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, base_y-Inches(0.12), Inches(0.28), Inches(0.28))
            dot.fill.solid(); dot.fill.fore_color.rgb = BLUE
            dot.line.fill.background()
            self.add_text(slide, item.get("time", f"T{i+1}"), x-Inches(0.35), base_y-Inches(0.65), Inches(1.0), Inches(0.25), 12, True, BLUE_DARK, PP_ALIGN.CENTER)
            self.add_card(slide, x-Inches(0.55), base_y+Inches(0.38), Inches(1.55), Inches(1.2), item.get("title", ""), item.get("body", ""), [BLUE, GREEN, ORANGE][i%3])

    def add_thanks(self):
        slide = self.blank(); self.add_top_bar(slide)
        self.add_text(slide, self.spec.get("thanks", "感谢聆听"), Inches(0.75), Inches(2.45), Inches(6.5), Inches(0.8), 36, True, BLUE_DARK)
        self.add_text(slide, "Q & A", Inches(0.78), Inches(3.35), Inches(3), Inches(0.45), 24, True, BLUE)
        self.add_card(slide, Inches(7.2), Inches(2.2), Inches(4.8), Inches(2.1), "交付检查", "结构清晰\n视觉统一\n内容可编辑\n适合答辩/汇报", GREEN)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", required=True)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()
    spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    deck = PPTMaster(spec).build()
    deck.save(out)
    report = out.with_suffix(".report.txt")
    report.write_text(f"PPT generated: {out}\nslides: {len(deck.slides)}\ntitle: {spec.get('title')}\n", encoding="utf-8")
    print(f"PPT generated: {out}")
    print(f"Report: {report}")


if __name__ == "__main__":
    main()
